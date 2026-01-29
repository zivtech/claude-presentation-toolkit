"""
Brand Compliance Analyzer

Scans a PPTX file and identifies slides that need manual intervention
for brand compliance (hard-coded fonts, off-brand colors).
"""

import sys
from pathlib import Path
from typing import Optional, List, Dict, Any, Union

from .config import BrandConfig
from .pptx_utils import rgb_to_hex


def analyze_shape(shape, config: BrandConfig) -> List[str]:
    """Analyze a shape for brand compliance issues.

    Args:
        shape: python-pptx Shape object
        config: Brand configuration

    Returns:
        List of issue descriptions
    """
    issues = []

    # Check text frames
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                # Check font
                if run.font.name and config.is_bad_font(run.font.name):
                    issues.append(f"Hard-coded font: {run.font.name}")

                # Check font color
                try:
                    if run.font.color.type is not None and run.font.color.rgb:
                        hex_color = rgb_to_hex(run.font.color.rgb)
                        if not config.is_brand_color(hex_color):
                            issues.append(f"Off-brand text color: #{hex_color}")
                except:
                    pass

    # Check fill colors
    if hasattr(shape, 'fill'):
        try:
            if shape.fill.type is not None:
                if hasattr(shape.fill, 'fore_color') and shape.fill.fore_color.rgb:
                    hex_color = rgb_to_hex(shape.fill.fore_color.rgb)
                    if not config.is_brand_color(hex_color):
                        issues.append(f"Off-brand fill color: #{hex_color}")
        except:
            pass

    # Check line colors
    if hasattr(shape, 'line'):
        try:
            if shape.line.color.type is not None and shape.line.color.rgb:
                hex_color = rgb_to_hex(shape.line.color.rgb)
                if not config.is_brand_color(hex_color):
                    issues.append(f"Off-brand line color: #{hex_color}")
        except:
            pass

    return list(set(issues))  # Remove duplicates


def get_slide_title(slide) -> str:
    """Extract slide title if present."""
    if slide.shapes.title:
        title = slide.shapes.title.text.strip()
        return title[:50] + "..." if len(title) > 50 else title

    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            text = shape.text_frame.text.strip()
            return text[:50] + "..." if len(text) > 50 else text
    return "(No title)"


def analyze_presentation(
    pptx_path: Union[str, Path],
    config: BrandConfig,
    verbose: bool = True
) -> List[Dict[str, Any]]:
    """Analyze entire presentation for brand compliance.

    Args:
        pptx_path: Path to PPTX file
        config: Brand configuration
        verbose: Print detailed output

    Returns:
        List of dicts with slide issues
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx is required. Install with: pip install python-pptx")

    if verbose:
        print(f"\nAnalyzing: {pptx_path}")
        print("=" * 70)

    prs = Presentation(pptx_path)

    slides_with_issues = []
    all_fonts = set()

    for i, slide in enumerate(prs.slides, 1):
        slide_issues = []

        for shape in slide.shapes:
            issues = analyze_shape(shape, config)
            slide_issues.extend(issues)

            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            all_fonts.add(run.font.name)

        if slide_issues:
            slides_with_issues.append({
                'num': i,
                'title': get_slide_title(slide),
                'issues': list(set(slide_issues))
            })

    # Categorize issues
    font_issues = []
    color_issues = []
    both_issues = []

    for slide in slides_with_issues:
        has_font = any('font' in i.lower() for i in slide['issues'])
        has_color = any('color' in i.lower() for i in slide['issues'])

        if has_font and has_color:
            both_issues.append(slide)
        elif has_font:
            font_issues.append(slide)
        elif has_color:
            color_issues.append(slide)

    if verbose:
        _print_analysis_results(
            prs, slides_with_issues, all_fonts, config,
            both_issues, font_issues, color_issues
        )

    return slides_with_issues


def _print_analysis_results(
    prs, slides_with_issues, all_fonts, config,
    both_issues, font_issues, color_issues
):
    """Print formatted analysis results."""
    print(f"\nTotal slides: {len(prs.slides)}")
    print(f"Clean slides: {len(prs.slides) - len(slides_with_issues)}")
    print(f"Slides needing fixes: {len(slides_with_issues)}")

    if all_fonts:
        bad_fonts_found = [f for f in all_fonts if config.is_bad_font(f)]
        if bad_fonts_found:
            print(f"\nNon-brand fonts found: {', '.join(sorted(bad_fonts_found))}")

    if both_issues:
        print(f"\n{'='*70}")
        print(f"HIGH PRIORITY - Font AND Color Issues ({len(both_issues)} slides)")
        print("-" * 50)
        for slide in both_issues:
            print(f"\nSlide {slide['num']}: {slide['title']}")
            for issue in sorted(set(slide['issues'])):
                print(f"  - {issue}")

    if font_issues:
        print(f"\n{'='*70}")
        print(f"MEDIUM PRIORITY - Font Issues ({len(font_issues)} slides)")
        print("-" * 50)
        for slide in font_issues:
            print(f"\nSlide {slide['num']}: {slide['title']}")
            for issue in sorted(set(slide['issues'])):
                print(f"  - {issue}")

    if color_issues:
        print(f"\n{'='*70}")
        print(f"LOWER PRIORITY - Color Issues ({len(color_issues)} slides)")
        print("-" * 50)
        for slide in color_issues:
            print(f"\nSlide {slide['num']}: {slide['title']}")
            for issue in sorted(set(slide['issues'])):
                print(f"  - {issue}")

    # Summary
    print(f"\n{'='*70}")
    print("SUMMARY")
    print("=" * 70)
    print(f"Total slides: {len(prs.slides)}")
    print(f"Clean (no issues): {len(prs.slides) - len(slides_with_issues)}")
    print(f"Need attention: {len(slides_with_issues)}")
    print(f"  - High priority (font+color): {len(both_issues)}")
    print(f"  - Medium priority (font): {len(font_issues)}")
    print(f"  - Lower priority (color): {len(color_issues)}")

    if slides_with_issues:
        print(f"\nSlides needing fixes: {', '.join(str(s['num']) for s in slides_with_issues)}")

    # Recommended actions
    print(f"\n{'='*70}")
    print("RECOMMENDED ACTIONS")
    print("=" * 70)

    bad_fonts_found = [f for f in all_fonts if config.is_bad_font(f)]
    if bad_fonts_found:
        print("\n1. BULK FONT REPLACEMENT (PowerPoint: Home -> Replace -> Replace Fonts)")
        for font in sorted(bad_fonts_found):
            # Suggest appropriate brand font
            brand_fonts = config.fonts.brand
            suggested = brand_fonts[0] if brand_fonts else "Brand Font"
            print(f"   {font} -> {suggested}")

    if color_issues or both_issues:
        print("\n2. FIX OFF-BRAND COLORS")
        print("   Review flagged slides and update to brand palette:")
        colors = config.colors.all_colors()
        color_list = list(colors.items())[:3]  # Show first 3
        print("   " + " | ".join(f"{name}: #{hex_val.upper()}" for name, hex_val in color_list))

    print("\n3. RUN THIS SCRIPT AGAIN to verify fixes")


def get_analysis_json(slides_with_issues: List[Dict[str, Any]], total_slides: int) -> Dict[str, Any]:
    """Convert analysis results to JSON-serializable format.

    Args:
        slides_with_issues: List of slide issue dicts
        total_slides: Total number of slides in presentation

    Returns:
        JSON-serializable dict
    """
    font_issues = []
    color_issues = []
    both_issues = []

    for slide in slides_with_issues:
        has_font = any('font' in i.lower() for i in slide['issues'])
        has_color = any('color' in i.lower() for i in slide['issues'])

        if has_font and has_color:
            both_issues.append(slide)
        elif has_font:
            font_issues.append(slide)
        elif has_color:
            color_issues.append(slide)

    return {
        'total_slides': total_slides,
        'clean_slides': total_slides - len(slides_with_issues),
        'slides_with_issues': len(slides_with_issues),
        'high_priority': len(both_issues),
        'medium_priority': len(font_issues),
        'low_priority': len(color_issues),
        'issues': slides_with_issues,
    }
